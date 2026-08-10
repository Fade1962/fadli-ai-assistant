from pptx import Presentation


def process_presentation(filename):

    try:

        presentation = Presentation(
            filename
        )

        output = []


        # ==========================================
        # SLIDES
        # ==========================================

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1
        ):

            output.append(
                f"=== SLIDE {slide_number} ==="
            )


            # ======================================
            # SHAPES / TEXT
            # ======================================

            for shape in slide.shapes:

                if hasattr(
                    shape,
                    "text"
                ):

                    text = shape.text.strip()

                    if text:

                        output.append(
                            text
                        )


                # ==================================
                # TABLE
                # ==================================

                if shape.has_table:

                    for row in shape.table.rows:

                        cells = []

                        for cell in row.cells:

                            value = (
                                cell.text
                                .strip()
                            )

                            cells.append(
                                value
                            )


                        if any(cells):

                            output.append(
                                " | ".join(cells)
                            )


        result = "\n\n".join(
            output
        ).strip()


        if not result:

            return (
                "PowerPoint berhasil dibaca, "
                "tetapi tidak ditemukan teks."
            )


        return result


    except Exception as error:

        print(
            "PRESENTATION PROCESS ERROR:",
            repr(error)
        )

        raise
