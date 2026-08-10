from pptx import Presentation
from pptx.util import Pt

def _chunks(lines, max_chars=1100):
    chunk, size = [], 0
    for line in lines:
        if chunk and size + len(line) > max_chars:
            yield chunk
            chunk, size = [], 0
        chunk.append(line)
        size += len(line) + 1
    if chunk:
        yield chunk

def create_pptx(filename, content):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Fadli AI"
    slide.placeholders[1].text = "Generated from Telegram request"

    lines = [x.strip() for x in (content or "").splitlines() if x.strip()]
    for i, block in enumerate(_chunks(lines), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Ringkasan {i}"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for j, line in enumerate(block):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
    prs.save(filename)
    return filename
